# -*- coding: utf-8 -*-
"""Rebuild the rendered layout as an A4 .pptx whose every element is editable
   (and which Canva imports as an editable design)."""
import json, os, re, sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from issues import resolve

BASE = os.path.dirname(os.path.abspath(__file__))
B = os.path.join(BASE, 'build')
KEY = sys.argv[1] if len(sys.argv) > 1 else '2026-07'
CFG = resolve(KEY)

layout = json.load(open(os.path.join(B, f'layout-{KEY}.json')))
PW_PX, PH_PX = layout[0]['w'], layout[0]['h']
A4_W_MM, A4_H_MM = 210.0, 297.0
EMU_MM = 36000
SX = A4_W_MM / PW_PX          # px -> mm
FONT = 'Mukta'

def emu(mm): return Emu(int(round(mm * EMU_MM)))
def px(v):   return v * SX                      # px -> mm
def pt_of(px_size): return px_size * 0.75       # css px -> pt

def rgb(css):
    m = re.match(r'rgba?\(([\d.]+),\s*([\d.]+),\s*([\d.]+)', css or '')
    if not m: return RGBColor(0, 0, 0)
    return RGBColor(*(int(round(float(m.group(i)))) for i in (1, 2, 3)))

ALIGN = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT,
         'left': PP_ALIGN.LEFT, 'start': PP_ALIGN.LEFT, 'justify': PP_ALIGN.JUSTIFY}

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'


def set_geom(shape, prst, adj=None):
    """Give a picture (or shape) a preset geometry — roundRect / ellipse —
       so CSS border-radius survives the export."""
    spPr = shape._element.spPr
    for tag in ('prstGeom', 'custGeom'):
        old = spPr.find(A + tag)
        if old is not None:
            spPr.remove(old)
    geom = spPr.makeelement(A + 'prstGeom', {'prst': prst})
    lst = spPr.makeelement(A + 'avLst', {})
    if adj is not None:
        gd = spPr.makeelement(A + 'gd', {'name': 'adj',
                                         'fmla': 'val %d' % int(adj)})
        lst.append(gd)
    geom.append(lst)
    # prstGeom must precede any fill/line children
    ref = None
    for ch in spPr:
        if ch.tag in (A + 'noFill', A + 'solidFill', A + 'blipFill',
                      A + 'gradFill', A + 'pattFill', A + 'grpFill', A + 'ln'):
            ref = ch
            break
    if ref is None:
        spPr.append(geom)
    else:
        spPr.insert(list(spPr).index(ref), geom)


def set_alpha(pic, opacity):
    """Apply CSS opacity to a picture via alphaModFix."""
    if opacity is None or opacity >= 0.999:
        return
    blip = pic._element.blipFill.find(A + 'blip')
    if blip is None:
        return
    fx = blip.makeelement(A + 'alphaModFix',
                          {'amt': str(int(round(opacity * 100000)))})
    blip.append(fx)


def set_edge(shape, color, width_mm):
    if not color or width_mm <= 0:
        return
    shape.line.color.rgb = rgb(color)
    shape.line.width = Emu(int(round(width_mm * EMU_MM)))

def set_script_fonts(run, name):
    """Point the east-asian and complex-script slots at the same face, so
    Devanagari is not silently swapped for a latin fallback.  DrawingML
    requires the order latin -> ea -> cs; anything else makes PowerPoint
    report the file as damaged."""
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(A + 'latin')
    if latin is None:
        latin = rPr.makeelement(A + 'latin', {'typeface': name})
        rPr.append(latin)
    at = list(rPr).index(latin) + 1
    for tag in ('ea', 'cs'):
        el = rPr.find(A + tag)
        if el is None:
            el = rPr.makeelement(A + tag, {'typeface': name})
            rPr.insert(at, el)
        else:
            el.set('typeface', name)
        at = list(rPr).index(el) + 1

prs = Presentation()
prs.slide_width  = emu(A4_W_MM)
prs.slide_height = emu(A4_H_MM)
blank = prs.slide_layouts[6]

TEAL = RGBColor(0x12, 0x61, 0x5B)

# PowerPoint sets the first baseline slightly lower than the browser for
# a given exact line spacing; measured against the rendered pages.
V_NUDGE = -0.34   # mm

for pg in layout:
    sl = prs.slides.add_slide(blank)

    # the footer rule is a CSS pseudo-element, so it is drawn here directly
    if pg['page'] > 1:
        bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  emu(19), emu(A4_H_MM - 11 - 0.92),
                                  emu(A4_W_MM - 38), emu(0.92))
        bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background(); bar.shadow.inherit = False

    for s in pg['shapes']:
        X, Y, W, H = px(s['x']), px(s['y']), px(s['w']), px(s['h'])
        if W <= 0.2 or H <= 0.2:      # a zero/negative extent corrupts the file
            continue

        if s['kind'] == 'rect':
            shape_t = MSO_SHAPE.OVAL if s.get('round') else (
                MSO_SHAPE.ROUNDED_RECTANGLE if s.get('radius', 0) > 1 else MSO_SHAPE.RECTANGLE)
            sh = sl.shapes.add_shape(shape_t, emu(X), emu(Y), emu(W), emu(H))
            sh.fill.solid(); sh.fill.fore_color.rgb = rgb(s['fill'])
            sh.line.fill.background(); sh.shadow.inherit = False
            if shape_t == MSO_SHAPE.ROUNDED_RECTANGLE:
                # adj is the corner radius as a fraction of the short side
                frac = px(s['radius']) / max(min(W, H), 0.01)
                set_geom(sh, 'roundRect', min(frac, 0.5) * 100000)
            set_edge(sh, s.get('edgeColor'), px(s.get('edgeW', 0)))

        elif s['kind'] == 'img':
            path = os.path.join(B, s['src'])
            if os.path.exists(path):
                pic = sl.shapes.add_picture(path, emu(X), emu(Y), emu(W), emu(H))
                if s.get('circle'):
                    set_geom(pic, 'ellipse')
                elif s.get('radius', 0) > 0.5:
                    frac = px(s['radius']) / max(min(W, H), 0.01)
                    set_geom(pic, 'roundRect', min(frac, 0.5) * 100000)
                set_edge(pic, s.get('edgeColor'), px(s.get('edgeW', 0)))
                set_alpha(pic, s.get('opacity'))
                if s.get('rot'):
                    pic.rotation = s['rot'] % 360

        elif s['kind'] == 'text':
            pad = 0.6
            if s.get('vert'):
                # swap the box about its centre, then turn it upright
                cx, cy = X + W / 2, Y + H / 2
                bw, bh = H, W
                tb = sl.shapes.add_textbox(emu(cx - bw / 2), emu(cy - bh / 2),
                                           emu(bw), emu(bh))
                tb.rotation = 270
            else:
                # keep the box exactly on the measured line band: padding the
                # left edge shifts every line, and widening asymmetrically
                # would move centred text off-centre
                slack = 2.0
                al = s.get('align', 'left')
                if al == 'center':
                    bx = X - slack / 2
                elif al == 'right':
                    bx = X - slack
                else:
                    bx = X
                tb = sl.shapes.add_textbox(emu(bx), emu(Y + V_NUDGE),
                                           emu(W + slack), emu(H))
            tf = tb.text_frame
            tf.word_wrap = not s.get('vert')   # a rotated line must not wrap
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = (PP_ALIGN.CENTER if s.get('vert')
                           else ALIGN.get(s.get('align', 'left'), PP_ALIGN.LEFT))
            size = pt_of(s['size'])
            p.line_spacing = Pt(pt_of(s['lh']))
            for r in s['runs']:
                run = p.add_run()
                run.text = r['t']
                f = run.font
                f.name = FONT
                f.size = Pt(round(size, 1))
                f.bold = bool(r.get('b'))
                f.italic = bool(r.get('i'))
                f.underline = bool(s.get('underline'))
                f.color.rgb = rgb(s['color'])
                if s.get('spacing'):
                    # spc is in hundredths of a point
                    run._r.get_or_add_rPr().set(
                        'spc', str(int(round(pt_of(s['spacing']) * 100))))
                set_script_fonts(run, FONT)

out = os.path.join(B, CFG['out'] + '_canva.pptx')
prs.save(out)
print('wrote', out, os.path.getsize(out) // 1024, 'KB;', len(layout), 'slides')
